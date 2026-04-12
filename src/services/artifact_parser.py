"""Artifact parser service: PPTX, PDF, GitHub README extraction + LLM structured extraction.

Parses project artifacts (presentations, PDFs, GitHub READMEs) and uses LLM
to extract structured ProjectExtraction data from raw text.
"""

import io
import json
import logging
import re

import httpx

logger = logging.getLogger(__name__)


async def parse_pptx(url_or_path: str) -> str:
    """Extract text from PPTX file. Handles URLs and local paths."""
    from pptx import Presentation

    if url_or_path.startswith("http"):
        async with httpx.AsyncClient(timeout=30) as client:
            resp = await client.get(url_or_path)
            resp.raise_for_status()
            data = io.BytesIO(resp.content)
    else:
        data = url_or_path

    prs = Presentation(data)
    slides_text: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        if texts:
            slides_text.append(f"Слайд {i}: {' | '.join(texts)}")

    return "\n".join(slides_text)


async def parse_pdf(url_or_path: str) -> str:
    """Extract text from PDF file."""
    import pymupdf

    if url_or_path.startswith("http"):
        async with httpx.AsyncClient(timeout=30) as client:
            resp = await client.get(url_or_path)
            resp.raise_for_status()
            data = resp.content
    else:
        with open(url_or_path, "rb") as f:
            data = f.read()

    doc = pymupdf.open(stream=data, filetype="pdf")
    pages_text: list[str] = []
    for i, page in enumerate(doc, 1):
        text = page.get_text().strip()
        if text:
            pages_text.append(f"Страница {i}: {text[:500]}")
    doc.close()

    return "\n".join(pages_text)


async def parse_github_readme(github_url: str) -> str:
    """Fetch README from GitHub repo."""
    match = re.match(r"https?://github\.com/([^/]+)/([^/]+)", github_url)
    if not match:
        return ""
    owner, repo = match.group(1), match.group(2).rstrip(".git")

    async with httpx.AsyncClient(timeout=15) as client:
        for branch in ["main", "master"]:
            url = f"https://raw.githubusercontent.com/{owner}/{repo}/{branch}/README.md"
            resp = await client.get(url)
            if resp.status_code == 200:
                return resp.text[:3000]
    return ""


async def parse_presentation(url: str) -> str:
    """Auto-detect format (PPTX or PDF) and parse."""
    url_lower = url.lower()
    if url_lower.endswith(".pptx"):
        return await parse_pptx(url)
    elif url_lower.endswith(".pdf"):
        return await parse_pdf(url)
    else:
        # Try PDF first (more common), fallback to PPTX
        try:
            return await parse_pdf(url)
        except Exception:
            try:
                return await parse_pptx(url)
            except Exception:
                logger.warning("Could not parse presentation: %s", url)
                return ""


async def extract_structured(
    raw_text: str,
    project_title: str,
    project_description: str,
    platform_client,
) -> dict:
    """Use LLM to extract structured ProjectExtraction from raw text."""
    from src.schemas.tools import ProjectExtraction

    prompt = (
        f"Проект: {project_title}\n"
        f"Описание: {project_description}\n\n"
        f"Текст из артефактов (презентация/README):\n{raw_text[:5000]}\n\n"
        "Извлеки информацию в ТОЧНО таком JSON формате:\n"
        "{\n"
        '  "problem": "какую проблему решает (1-2 предложения)",\n'
        '  "solution": "как решает (1-2 предложения)",\n'
        '  "audience": "для кого",\n'
        '  "stack": ["Python", "PyTorch"],\n'
        '  "novelty": "в чем новизна",\n'
        '  "risks": "ограничения или null",\n'
        '  "key_metrics": ["accuracy 94%", "F1=0.91", "1000 rps"] или null,\n'
        '  "production_readiness": "prototype" или "mvp" или "production" или null,\n'
        '  "team_size": число или null,\n'
        '  "red_flags": [{"category": "metric", "description": "текст", "severity": "low"}] или null\n'
        "}\n\n"
        "ПРАВИЛА:\n"
        "- Поля problem, solution, audience, stack, novelty ОБЯЗАТЕЛЬНЫ\n"
        "- key_metrics - СПИСОК СТРОК, не словарь\n"
        "- Извлекай ТОЛЬКО то, что ЯВНО указано в тексте\n"
        "- Если данных нет - null (кроме обязательных полей)"
    )

    system = (
        "Ты аналитик студенческих AI-проектов. "
        "Извлекай факты строго в указанном JSON формате. "
        "НЕ выдумывай данные. НЕ добавляй свои поля. НЕ меняй имена полей."
    )

    try:
        resp = await platform_client.chat_completion(
            messages=[
                {"role": "system", "content": system},
                {"role": "user", "content": prompt},
            ],
            response_format={"type": "json_object"},
        )
        content = resp["choices"][0]["message"]["content"]
        data = json.loads(content)

        # Fix common LLM mistakes before validation
        if isinstance(data.get("key_metrics"), dict):
            # LLM returned dict instead of list - convert to list of strings
            data["key_metrics"] = [
                f"{k}: {v}" for k, v in data["key_metrics"].items() if v is not None
            ]
        if data.get("stack") is None:
            data["stack"] = []

        try:
            extraction = ProjectExtraction(**data)
            return extraction.model_dump()
        except Exception:
            # Validation failed - store raw data as-is (JSONB accepts any dict)
            logger.warning("Validation partial for %s, storing raw", project_title)
            # Ensure minimum required fields
            data.setdefault("problem", project_description[:100])
            data.setdefault("solution", "Не извлечено")
            data.setdefault("audience", "Не указано")
            data.setdefault("stack", [])
            data.setdefault("novelty", "Не извлечено")
            return data
    except Exception as e:
        logger.error("Structured extraction failed for %s: %s", project_title, e)
        return {}
